import urllib
import requests
from bs4 import BeautifulSoup
from time import sleep
import os
import sys
import argparse
import pptx
import imageio
import glob
import pptx.util
import pptx
import json

HEADERS = {'User-Agent': 'Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/51.0.2704.106 Safari/537.36'}


def get_webpage_links(url, base_url):
    request = requests.get(url, HEADERS)
    soup = BeautifulSoup(request.content, 'html.parser')

    urls = []

    for ab in soup.find_all(class_='details text-left'):
        name = ab.find('a')['href']
        urls.append((base_url[:-1]+name, name.split('/')[2].split('?')[0]))

    
    return urls


def download_images(output_folder, HEADERS, url, name):
    request = requests.get(url, HEADERS)
    soup = BeautifulSoup(request.content, 'html.parser')
    images = soup.findAll('img', {'class':'slide_image'})

    dd = dict()
    meta_data = dict()

    #namee = soup.find(class_='details text-left')
    meta_data['name'] = name.replace('-',' ')
    extract = soup.find('div', {'class':'author-text'})
    meta_data['author'] = extract.find('a').find('span').get_text()
    if(extract.find('small') is not None):
        meta_data['author'] += extract.find('small').get_text()

    meta_data['year'] = 'N/A'
    meta_data['source'] = 'LinkedIn Slideshare'
    dd[meta_data['name']] = meta_data

    i = 0
    for image in images:
        try:
            image_url = image.get('data-full').split('?')[0]
            urllib.request.urlretrieve(image_url, os.path.join(output_folder, 'ab'+str(i)+'.png'))
        except Exception as e:
            print('Failed to download because of', e)
        i += 1
    return dd


def convert_images_to_ppt(output_folder, images_folder, file_name):
    OUTPUT_TAG = "PRESENTATION"

    # new
    prs = pptx.Presentation()
    prs.slide_height = 5143500

    pic_left  = int(prs.slide_width * 0.15)
    pic_top   = int(prs.slide_height * 0.1)
    pic_width = int(prs.slide_width * 0.7)


    for g in os.listdir(images_folder):
        g = os.path.join(images_folder, g)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        img = imageio.imread(g)
        pic_height = int(pic_width * img.shape[0] / img.shape[1])
        pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
    prs.save(os.path.join(output_folder, file_name))


def delete_images(output_folder):
    for file in os.listdir(output_folder):
        if(file.endswith('png')):
            os.remove(os.path.join(output_folder, file))

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("url", help="Topic to retrieve all the ppts.")
    args = parser.parse_args()

    topic = '+'.join(args.url.lower().split(' '))#+'+topic'

    meta_data = dict()
    base_url = 'https://www.slideshare.net/'
    output_folder = os.path.join(os.getcwd(), 'PPTs')
    images_folder = os.path.join(os.getcwd(), 'pptimages')
    
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    if not os.path.exists(images_folder):
        os.makedirs(images_folder)

    for i in range(1, 2):
        if(i == 0):
            url = 'https://www.slideshare.net/search/slideshow?searchfrom=header&q='+topic
        else:
            url = 'https://www.slideshare.net/search/slideshow?lang=%2A%2A&page='+str(i)+'&q='+topic+'&searchfrom=header&sort=relevance'
        urls = get_webpage_links(url, base_url)
        
        v = 1
        for each in urls:
            print(each[0])
            meta_d = download_images(images_folder, HEADERS, each[0], each[1])
            meta_data = dict(meta_data, **meta_d)
            convert_images_to_ppt(output_folder, images_folder, each[1]+'.pptx')
            delete_images(images_folder)
            print('\n')
            v += 1

    return_dic = {'meta': meta_data}
    with open('metadata_ppt.json', 'w') as file:
        file.write(json.dumps(return_dic))

if __name__ == "__main__":
    main()
